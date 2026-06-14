#!/usr/bin/env python3
"""Merge generated slide images into PPTX and PDF.

Usage:
    python3 merge_deck.py /path/to/paper-deck/topic-slug
"""

from __future__ import annotations

import argparse
import re
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches


IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp"}
SLIDE_W_IN = 13.333333
SLIDE_H_IN = 7.5


def natural_key(path: Path) -> tuple[int, str]:
    match = re.match(r"^(\d+)", path.stem)
    number = int(match.group(1)) if match else 9999
    return number, path.name


def find_images(deck_dir: Path) -> list[Path]:
    image_dir = deck_dir / "images"
    if not image_dir.exists():
        raise SystemExit(f"Missing images directory: {image_dir}")

    images = [p for p in image_dir.iterdir() if p.suffix.lower() in IMAGE_EXTS]
    images.sort(key=natural_key)
    if not images:
        raise SystemExit(f"No slide images found in: {image_dir}")
    return images


def make_pptx(images: list[Path], output: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank = prs.slide_layouts[6]

    for image in images:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(image),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    prs.save(output)


def make_pdf(images: list[Path], output: Path) -> None:
    frames: list[Image.Image] = []
    for image_path in images:
        with Image.open(image_path) as img:
            rgb = img.convert("RGB")
            frames.append(rgb.copy())

    first, *rest = frames
    first.save(output, save_all=True, append_images=rest, resolution=150.0)


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("deck_dir", help="Deck directory containing images/")
    parser.add_argument("--name", help="Output base name. Defaults to deck directory name.")
    args = parser.parse_args()

    deck_dir = Path(args.deck_dir).expanduser().resolve()
    if not deck_dir.exists():
        raise SystemExit(f"Deck directory does not exist: {deck_dir}")

    images = find_images(deck_dir)
    base = args.name or deck_dir.name
    pptx_path = deck_dir / f"{base}.pptx"
    pdf_path = deck_dir / f"{base}.pdf"

    make_pptx(images, pptx_path)
    make_pdf(images, pdf_path)

    print(f"Merged {len(images)} slides")
    print(f"PPTX: {pptx_path}")
    print(f"PDF:  {pdf_path}")


if __name__ == "__main__":
    main()
